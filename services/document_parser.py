"""
Document Parser – extracts text from various file formats.

Supported: PDF, DOCX, PPTX, XLSX, CSV, TXT, JSON
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

MAX_FALLBACK_CHARS = 300_000
MAX_EXTRACTED_CHARS = 200_000
MAX_PDF_PAGES = 120
MAX_DOCX_PARAGRAPHS = 8000
MAX_PPTX_SLIDES = 500
MAX_XLSX_ROWS = 20000
MAX_CSV_ROWS = 20000


class DocumentParser:
    """Stateless document text extractor."""

    SUPPORTED_TYPES = {"pdf", "docx", "pptx", "xlsx", "csv", "txt", "json"}

    # ── Public API ───────────────────────────────────────
    def parse(self, file_bytes: bytes, file_type: str) -> str:
        """
        Extract plain text from raw bytes of a document.

        Parameters
        ----------
        file_bytes : bytes
            Raw file content.
        file_type : str
            Lowercase extension (e.g. ``"pdf"``).

        Returns
        -------
        str
            Extracted text.
        """
        file_type = file_type.lower().lstrip(".")
        if file_type not in self.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {file_type}")

        handler = getattr(self, f"_parse_{file_type}")
        text: str = handler(file_bytes)
        text = self._normalize_text(text, MAX_EXTRACTED_CHARS)
        logger.info("Parsed %s document – %d chars extracted", file_type, len(text))
        return text

    def extract_metadata(self, file_bytes: bytes, file_type: str) -> dict[str, Any]:
        """Extract basic metadata from a document."""
        file_type = file_type.lower().lstrip(".")
        return {
            "file_type": file_type,
            "size_bytes": len(file_bytes),
            "char_count": len(self.parse(file_bytes, file_type)),
        }

    # ── Private parsers ──────────────────────────────────
    @staticmethod
    def _best_effort_text(data: bytes) -> str:
        text = data.decode("utf-8", errors="ignore")
        if not text.strip():
            text = data.decode("latin-1", errors="ignore")
        text = DocumentParser._normalize_text(text, MAX_FALLBACK_CHARS)
        if len(text) > MAX_FALLBACK_CHARS:
            logger.warning(
                "Fallback parser produced very large text (%d chars). Truncating to %d.",
                len(text),
                MAX_FALLBACK_CHARS,
            )
            text = text[:MAX_FALLBACK_CHARS]
        return text or "Document content could not be parsed, but file upload succeeded."

    @staticmethod
    def _normalize_text(text: str, max_chars: int) -> str:
        # Keep layout separators, remove binary noise, collapse excessive whitespace.
        text = re.sub(r"[^\x09\x0A\x0D\x20-\x7E]", " ", text)
        text = re.sub(r"[ \t]{2,}", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = text.strip()
        if len(text) > max_chars:
            text = text[:max_chars]
        return text

    @staticmethod
    def _parse_pdf(data: bytes) -> str:
        import fitz  # PyMuPDF

        text_parts: list[str] = []
        extracted = 0
        try:
            with fitz.open(stream=data, filetype="pdf") as doc:
                for page_index, page in enumerate(doc):
                    if page_index >= MAX_PDF_PAGES:
                        logger.warning("PDF page cap reached (%d). Stopping extraction.", MAX_PDF_PAGES)
                        break
                    page_text = page.get_text() or ""
                    if page_text:
                        text_parts.append(page_text)
                        extracted += len(page_text)
                    if extracted >= MAX_EXTRACTED_CHARS:
                        logger.warning(
                            "PDF text cap reached (%d chars). Stopping extraction.",
                            MAX_EXTRACTED_CHARS,
                        )
                        break
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning("PDF structured parsing failed, using text fallback: %s", e)
            return DocumentParser._best_effort_text(data)

    @staticmethod
    def _parse_docx(data: bytes) -> str:
        from docx import Document as DocxDocument

        try:
            doc = DocxDocument(io.BytesIO(data))
            text_parts: list[str] = []
            extracted = 0
            for idx, para in enumerate(doc.paragraphs):
                if idx >= MAX_DOCX_PARAGRAPHS:
                    logger.warning("DOCX paragraph cap reached (%d).", MAX_DOCX_PARAGRAPHS)
                    break
                para_text = (para.text or "").strip()
                if not para_text:
                    continue
                text_parts.append(para_text)
                extracted += len(para_text)
                if extracted >= MAX_EXTRACTED_CHARS:
                    break
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning("DOCX structured parsing failed, using text fallback: %s", e)
            return DocumentParser._best_effort_text(data)

    @staticmethod
    def _parse_pptx(data: bytes) -> str:
        from pptx import Presentation

        try:
            prs = Presentation(io.BytesIO(data))
            text_parts: list[str] = []
            extracted = 0
            for idx, slide in enumerate(prs.slides):
                if idx >= MAX_PPTX_SLIDES:
                    logger.warning("PPTX slide cap reached (%d).", MAX_PPTX_SLIDES)
                    break
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text = shape.text_frame.text or ""
                        if slide_text:
                            text_parts.append(slide_text)
                            extracted += len(slide_text)
                    if extracted >= MAX_EXTRACTED_CHARS:
                        break
                if extracted >= MAX_EXTRACTED_CHARS:
                    break
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning("PPTX structured parsing failed, using text fallback: %s", e)
            return DocumentParser._best_effort_text(data)

    @staticmethod
    def _parse_xlsx(data: bytes) -> str:
        from openpyxl import load_workbook

        try:
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            text_parts: list[str] = []
            extracted = 0
            row_count = 0
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_count += 1
                    if row_count > MAX_XLSX_ROWS:
                        logger.warning("XLSX row cap reached (%d).", MAX_XLSX_ROWS)
                        break
                    cells = [str(c) if c is not None else "" for c in row]
                    row_text = "\t".join(cells)
                    text_parts.append(row_text)
                    extracted += len(row_text)
                    if extracted >= MAX_EXTRACTED_CHARS:
                        break
                if row_count > MAX_XLSX_ROWS or extracted >= MAX_EXTRACTED_CHARS:
                    break
            return "\n".join(text_parts)
        except Exception as e:
            # Some files are mislabeled/corrupted .xlsx archives. Fallback to best-effort text
            # extraction so ingestion can continue rather than failing the whole document.
            logger.warning("XLSX structured parsing failed, using text fallback: %s", e)
            return DocumentParser._best_effort_text(data)

    @staticmethod
    def _parse_csv(data: bytes) -> str:
        text = data.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows: list[str] = []
        extracted = 0
        for idx, row in enumerate(reader):
            if idx >= MAX_CSV_ROWS:
                logger.warning("CSV row cap reached (%d).", MAX_CSV_ROWS)
                break
            row_text = "\t".join(row)
            rows.append(row_text)
            extracted += len(row_text)
            if extracted >= MAX_EXTRACTED_CHARS:
                break
        return "\n".join(rows)

    @staticmethod
    def _parse_txt(data: bytes) -> str:
        return data.decode("utf-8", errors="replace")[:MAX_EXTRACTED_CHARS]

    @staticmethod
    def _parse_json(data: bytes) -> str:
        obj = json.loads(data.decode("utf-8", errors="replace"))
        return json.dumps(obj, indent=2, ensure_ascii=False)[:MAX_EXTRACTED_CHARS]
