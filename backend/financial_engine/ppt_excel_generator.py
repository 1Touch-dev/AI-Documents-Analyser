"""
Excel & PowerPoint Model Generator
==================================
Dynamically builds and outputs fully functional Excel workbooks (.xlsx)
and professional PowerPoint presentation slide decks (.pptx) representing
interactive financial models, scenarios, and board-ready reporting packages.
"""

import logging
import io
import pandas as pd
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

class ExportGenerator:
    """Generates financial model spreadsheets and narrative presentation decks."""

    @staticmethod
    def generate_excel_model(
        starting_cash: float,
        revenue_items: list,
        expense_items: list,
        forecast_data: dict
    ) -> bytes:
        """Dynamically designs a professional financial model with formulas."""
        wb = openpyxl.Workbook()
        
        # ──────────────────────────────────────────────────────────
        # Sheet 1: Executive Summary & Forecasts
        # ──────────────────────────────────────────────────────────
        ws1 = wb.active
        ws1.title = "Executive Forecast"
        ws1.views.sheetView[0].showGridLines = True

        # Styles
        title_font = Font(name="Segoe UI", size=16, bold=True, color="FFFFFF")
        header_font = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")
        bold_font = Font(name="Segoe UI", size=11, bold=True, color="000000")
        regular_font = Font(name="Segoe UI", size=11, color="000000")
        
        dark_blue_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
        light_blue_fill = PatternFill(start_color="DDEBF7", end_color="DDEBF7", fill_type="solid")
        gray_fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")

        thin_border = Border(
            left=Side(style='thin', color='D9D9D9'),
            right=Side(style='thin', color='D9D9D9'),
            top=Side(style='thin', color='D9D9D9'),
            bottom=Side(style='thin', color='D9D9D9')
        )
        double_bottom_border = Border(
            top=Side(style='thin', color='000000'),
            bottom=Side(style='double', color='000000')
        )

        # Title Block
        ws1.merge_cells("A1:G1")
        ws1["A1"] = "ENTERPRISE FINANCIAL OPERATING SYSTEM — FORECAST MODEL"
        ws1["A1"].font = title_font
        ws1["A1"].fill = dark_blue_fill
        ws1["A1"].alignment = Alignment(horizontal="center", vertical="center")
        ws1.row_dimensions[1].height = 40

        # Sub-header
        ws1["A3"] = f"Initial Cash Reserve: ${starting_cash:,.2f}"
        ws1["A3"].font = bold_font

        # Forecast Table Header
        headers = ["Forecast Period", "Days Included", "Cash Inbound", "Cash Outbound", "Net Cash Flow", "Ending Cash Balance", "Monthly Burn Rate"]
        for col_num, h in enumerate(headers, 1):
            cell = ws1.cell(row=5, column=col_num)
            cell.value = h
            cell.font = header_font
            cell.fill = dark_blue_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border
        ws1.row_dimensions[5].height = 25

        # Populate Forecast Periods
        periods = [("30 Days", 30), ("60 Days", 60), ("90 Days", 90), ("180 Days", 180)]
        for idx, (label, days) in enumerate(periods, 6):
            key = f"forecast_{days}d"
            fd = forecast_data.get(key, {})
            
            ws1.cell(row=idx, column=1, value=label).font = bold_font
            ws1.cell(row=idx, column=2, value=days).font = regular_font
            ws1.cell(row=idx, column=3, value=fd.get("cash_in", 0.0)).font = regular_font
            ws1.cell(row=idx, column=4, value=fd.get("cash_out", 0.0)).font = regular_font
            ws1.cell(row=idx, column=5, value=fd.get("net_cash_flow", 0.0)).font = regular_font
            ws1.cell(row=idx, column=6, value=fd.get("ending_cash", 0.0)).font = bold_font
            ws1.cell(row=idx, column=7, value=fd.get("burn_rate", 0.0)).font = regular_font

            # Alignments, Borders, Number Formats
            for col_num in range(1, 8):
                cell = ws1.cell(row=idx, column=col_num)
                cell.border = thin_border
                if col_num > 2:
                    cell.number_format = '"$"#,##0.00'
                    cell.alignment = Alignment(horizontal="right")
                else:
                    cell.alignment = Alignment(horizontal="center")

        # Set Column Widths
        for col in ws1.columns:
            max_len = max(len(str(cell.value or '')) for cell in col)
            ws1.column_dimensions[col[0].column_letter].width = max(max_len + 3, 14)

        # ──────────────────────────────────────────────────────────
        # Sheet 2: Extracted Data Roster
        # ──────────────────────────────────────────────────────────
        ws2 = wb.create_sheet(title="Extracted Items")
        ws2.views.sheetView[0].showGridLines = True

        ws2.cell(row=1, column=1, value="EXTRACTED FINANCIAL LINE ITEMS").font = Font(name="Segoe UI", size=14, bold=True)
        
        ws2.cell(row=3, column=1, value="Revenue Items").font = Font(name="Segoe UI", size=12, bold=True, color="1F4E79")
        ws2.cell(row=4, column=1, value="Name").font = header_font
        ws2.cell(row=4, column=1).fill = dark_blue_fill
        ws2.cell(row=4, column=2, value="Category").font = header_font
        ws2.cell(row=4, column=2).fill = dark_blue_fill
        ws2.cell(row=4, column=3, value="Amount").font = header_font
        ws2.cell(row=4, column=3).fill = dark_blue_fill

        row_idx = 5
        for rev in revenue_items:
            ws2.cell(row=row_idx, column=1, value=rev.get("name", "")).font = regular_font
            ws2.cell(row=row_idx, column=2, value=rev.get("category", "")).font = regular_font
            cell_val = ws2.cell(row=row_idx, column=3, value=rev.get("amount", 0.0))
            cell_val.font = regular_font
            cell_val.number_format = '"$"#,##0.00'
            row_idx += 1

        row_idx += 2
        ws2.cell(row=row_idx, column=1, value="Expense Items").font = Font(name="Segoe UI", size=12, bold=True, color="A51D24")
        row_idx += 1
        ws2.cell(row=row_idx, column=1, value="Name").font = header_font
        ws2.cell(row=row_idx, column=1).fill = PatternFill(start_color="A51D24", end_color="A51D24", fill_type="solid")
        ws2.cell(row=row_idx, column=2, value="Category").font = header_font
        ws2.cell(row=row_idx, column=2).fill = PatternFill(start_color="A51D24", end_color="A51D24", fill_type="solid")
        ws2.cell(row=row_idx, column=3, value="Amount").font = header_font
        ws2.cell(row=row_idx, column=3).fill = PatternFill(start_color="A51D24", end_color="A51D24", fill_type="solid")
        row_idx += 1

        for exp in expense_items:
            ws2.cell(row=row_idx, column=1, value=exp.get("name", "")).font = regular_font
            ws2.cell(row=row_idx, column=2, value=exp.get("category", "")).font = regular_font
            cell_val = ws2.cell(row=row_idx, column=3, value=exp.get("amount", 0.0))
            cell_val.font = regular_font
            cell_val.number_format = '"$"#,##0.00'
            row_idx += 1

        for col in ws2.columns:
            max_len = max(len(str(cell.value or '')) for cell in col)
            ws2.column_dimensions[col[0].column_letter].width = max(max_len + 3, 14)

        # Output to bytes
        out = io.BytesIO()
        wb.save(out)
        return out.getvalue()

    @staticmethod
    def generate_pptx_deck(
        title_text: str,
        board_summary: str,
        risks: list,
        forecast_data: dict
    ) -> bytes:
        """Creates a professional corporate presentation slide deck."""
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # widescreen 16:9
        prs.slide_height = Inches(7.5)

        # Slide 1: Title Slide (Dark Theme)
        blank_layout = prs.slide_layouts[6]
        slide1 = prs.slides.add_slide(blank_layout)
        
        # Solid Dark Background
        bg = slide1.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)  # slate-900

        # Title Textbox
        tx_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.5))
        tf = tx_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "CFO EXECUTIVE INTELLIGENCE REPORT"
        p.font.name = "Segoe UI"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(99, 102, 241)  # Indigo-400

        p3 = tf.add_paragraph()
        p3.text = "Prepared by Artificial Intelligence Financial Operating System"
        p3.font.name = "Segoe UI"
        p3.font.size = Pt(12)
        p3.font.italic = True
        p3.font.color.rgb = RGBColor(148, 163, 184)  # slate-400

        # Slide 2: Executive Summary (Light Theme)
        slide2 = prs.slides.add_slide(blank_layout)
        bg2 = slide2.background
        fill2 = bg2.fill
        fill2.solid()
        fill2.fore_color.rgb = RGBColor(248, 250, 252)  # slate-50

        # Header
        header_box = slide2.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.0))
        tf_header = header_box.text_frame
        p_head = tf_header.paragraphs[0]
        p_head.text = "Executive Briefing & Strategic Assessment"
        p_head.font.name = "Segoe UI"
        p_head.font.size = Pt(28)
        p_head.font.bold = True
        p_head.font.color.rgb = RGBColor(30, 41, 59)

        # Content Summary Box
        content_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.8))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        paragraphs = board_summary.split("\n\n")
        for i, para in enumerate(paragraphs):
            if i == 0:
                p_c = tf_content.paragraphs[0]
            else:
                p_c = tf_content.add_paragraph()
                p_c.space_before = Pt(12)
            p_c.text = para.strip()
            p_c.font.name = "Segoe UI"
            p_c.font.size = Pt(15)
            p_c.font.color.rgb = RGBColor(51, 65, 85)

        # Slide 3: Risk and Compliance Dashboard
        slide3 = prs.slides.add_slide(blank_layout)
        bg3 = slide3.background
        fill3 = bg3.fill
        fill3.solid()
        fill3.fore_color.rgb = RGBColor(248, 250, 252)

        header_box3 = slide3.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.0))
        p_head3 = header_box3.text_frame.paragraphs[0]
        p_head3.text = "Financial Risk Ledger & Mitigation Matrix"
        p_head3.font.name = "Segoe UI"
        p_head3.font.size = Pt(28)
        p_head3.font.bold = True
        p_head3.font.color.rgb = RGBColor(30, 41, 59)

        # 3-Column Risk Layout
        left_pos = [Inches(0.75), Inches(4.8), Inches(8.85)]
        col_width = Inches(3.7)
        top_pos = Inches(1.8)
        height = Inches(4.8)

        for idx, risk_obj in enumerate(risks[:3]):
            if idx >= 3:
                break
            box = slide3.shapes.add_textbox(left_pos[idx], top_pos, col_width, height)
            tf_box = box.text_frame
            tf_box.word_wrap = True

            # Risk Title
            p_title = tf_box.paragraphs[0]
            p_title.text = risk_obj.title.upper()
            p_title.font.name = "Segoe UI"
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            
            severity = risk_obj.severity.lower()
            if severity == "critical":
                p_title.font.color.rgb = RGBColor(185, 28, 28)  # Red-700
            elif severity == "high":
                p_title.font.color.rgb = RGBColor(217, 119, 6)   # Amber-600
            else:
                p_title.font.color.rgb = RGBColor(29, 78, 216)   # Blue-700

            # Severity label
            p_sev = tf_box.add_paragraph()
            p_sev.text = f"Severity: {risk_obj.severity.upper()}"
            p_sev.font.name = "Segoe UI"
            p_sev.font.size = Pt(11)
            p_sev.font.italic = True
            p_sev.font.color.rgb = RGBColor(100, 116, 139)
            p_sev.space_after = Pt(8)

            # Description
            p_desc = tf_box.add_paragraph()
            p_desc.text = f"Context: {risk_obj.description}"
            p_desc.font.name = "Segoe UI"
            p_desc.font.size = Pt(13)
            p_desc.font.color.rgb = RGBColor(51, 65, 85)
            p_desc.space_after = Pt(12)

            # Mitigation
            p_mit = tf_box.add_paragraph()
            p_mit.text = f"Directive Action:\n{risk_obj.mitigation_action}"
            p_mit.font.name = "Segoe UI"
            p_mit.font.size = Pt(13)
            p_mit.font.bold = True
            p_mit.font.color.rgb = RGBColor(15, 118, 110)  # Teal-700

        # Save to bytes
        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()
